import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path.cwd()
OUT_DIR = ROOT / "content-research"


def load_snapshot(path):
    source = path.read_text(encoding="utf-8")
    marker = "window.__SITE_SNAPSHOT__="
    start = source.index(marker) + len(marker)
    end = source.index("const $ =")
    data = json.loads(source[start:end].strip().rstrip(";"))
    return source, start, end, data


def save_snapshot(path, source, start, end, data):
    next_source = source[:start] + json.dumps(data, ensure_ascii=False, separators=(",", ":")) + ";\n" + source[end:]
    path.write_text(next_source, encoding="utf-8")


slides = [
    {
        "title": "湖北楚云数航科技有限公司",
        "subtitle": "AI驱动的跨境电商、AI内容与算力园区运营平台",
        "bullets": [
            "公司围绕AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售四大业务展开。",
            "以高校人才、政府园区、跨境渠道、内容生产、算力资源和产业服务商为核心资源，形成可招商、可交付、可复制的产业平台。",
            "战略合作方向包含光谷超算中心等算力资源，支撑AI内容、电商运营、虚拟商品和企业智能化应用。"
        ],
    },
    {
        "title": "我们是一家什么公司",
        "subtitle": "不是单一培训机构，而是AI产业应用与OPC孵化的综合运营商",
        "bullets": [
            "面向高校、政府园区、创业团队、投资人、跨境商家、AI内容团队和企业客户，提供课程培训、项目孵化、系统工具、园区运营和资源撮合。",
            "以“平台统一管控+项目公司/店群实战+AI内容生产+算力资源支撑+园区载体落地”为核心模式。",
            "目标是把AI内容、跨境电商、创业孵化和算力应用从零散项目变成规模化、标准化、可监管的产业服务。"
        ],
    },
    {
        "title": "企业愿景",
        "subtitle": "让AI能力、跨境渠道与园区载体成为年轻人和中小企业的增长基础设施",
        "bullets": [
            "愿景：打造湖北领先、面向全国复制的AI数字贸易与OPC全域孵化平台。",
            "使命：让学生、创业者和企业用一套平台完成学习、内容生产、开店运营、算力使用、资源对接和项目成长。",
            "价值主张：政府看得见产业成果，高校拿得到实训转化，企业接得住AI应用，投资人看得到真实数据。"
        ],
    },
    {
        "title": "四大业务版块",
        "subtitle": "从内容、店群、园区到算力，构建AI产业应用闭环",
        "bullets": [
            "AI漫剧制作与培训：面向短剧、漫剧、虚拟人、AI视频和内容电商，提供课程、工具、剧本、制作流程和商业化训练。",
            "跨境店群投资与孵化：围绕TikTok、Ozon、速卖通、独立站等赛道，组织店群投资、运营陪跑、分润结算和项目孵化。",
            "OPC园区基建与运营：建设跨境产品展示、实景直播、创业路演、AI算力实训、工商托管和政企展示一体化载体。",
            "AI算力运营及销售：联合光谷超算中心等资源，开展算力采购、API中转、企业应用套壳、模型服务和算力衍生产品销售。"
        ],
    },
    {
        "title": "AI漫剧制作与培训",
        "subtitle": "把AI内容生产变成可学习、可交付、可商业化的训练体系",
        "bullets": [
            "课程内容覆盖AI剧本策划、分镜生成、角色一致性、配音字幕、短视频剪辑、账号矩阵和内容投流。",
            "适合高校实训、就业创业课程、园区创作者训练营、MCN/内容团队升级和企业品牌内容生产。",
            "商业路径包括培训收费、项目制作、账号运营、内容带货、品牌宣传片和AI漫剧IP孵化。"
        ],
    },
    {
        "title": "跨境店群投资与孵化",
        "subtitle": "以真实店铺和经营数据筛选项目，而不是只停留在课程培训",
        "bullets": [
            "通过店群管理系统统一承接选品、刊登、订单、库存、客服、复盘和风控，形成可追踪经营数据。",
            "投资人和企业可参与店铺托管、品牌出海、货源合作、渠道分成、项目投资和运营合伙。",
            "学生和创业者从实训进入真实店铺运营，用订单、GMV、利润、复购和成长档案证明项目价值。"
        ],
    },
    {
        "title": "OPC园区基建与运营",
        "subtitle": "让园区具备招商、孵化、展示、直播、培训和数据监管能力",
        "bullets": [
            "室内空间承接OPC工位、AI算力实训室、跨境店群运营区、校企洽谈室和项目路演区。",
            "户外及公共空间承接跨境产品展销、实景直播、创业活动、品牌路演和政府参观展示。",
            "园区运营输出标准化SOP、招商材料、分中心机制、政策申报材料和产业数据看板。"
        ],
    },
    {
        "title": "AI算力运营及销售",
        "subtitle": "以光谷超算中心战略合作为资源基础，服务AI内容与企业应用",
        "bullets": [
            "算力业务包括算力资源销售、API平台中转、模型调用服务、企业应用套壳和算力衍生产品运营。",
            "优先服务AI漫剧、跨境电商选品与客服、数字人直播、虚拟商品、企业知识库和园区数据大屏。",
            "通过与光谷超算中心等资源合作，形成“算力资源+应用场景+销售渠道+项目孵化”的商业闭环。"
        ],
    },
    {
        "title": "四类平台能力",
        "subtitle": "业务不是散点，而由平台能力持续沉淀",
        "bullets": [
            "AI内容与漫剧制作中台：沉淀剧本、角色、分镜、素材、模板、制作任务和商业发布流程。",
            "跨境店群运营管理平台：管理店铺、订单、库存、客服、财务、分润、投资和经营复盘。",
            "OPC园区数字化运营平台：管理工位、学员、项目、活动、政策申报、数据看板和政企展示。",
            "AI算力运营与销售平台：承接算力套餐、API调用、客户订单、模型应用、成本核算和渠道销售。"
        ],
    },
    {
        "title": "资源渠道与合作方",
        "subtitle": "高校、政府园区、跨境渠道、算力资源与产业服务商共同构成生态",
        "bullets": [
            "高校资源：共建数字贸易、AI内容制作、SYB/双创课程、实习实训基地和创新创业学分体系。",
            "政府与园区资源：围绕光谷科学岛、产业园区、创业载体、政策申报、赛事活动和招商展示展开合作。",
            "跨境与服务商资源：对接TikTok、Ozon、速卖通、独立站、ERP、货源池、支付物流和内容电商团队。",
            "算力资源：与光谷超算中心等资源建立战略合作，支撑AI内容、电商智能化和企业AI应用落地。"
        ],
    },
    {
        "title": "合作模式",
        "subtitle": "让每类伙伴都能清楚知道如何进入楚云数航生态",
        "bullets": [
            "高校合作：共建AI漫剧与跨境电商实训基地，输出课程、导师、项目、平台和就业创业成果。",
            "园区合作：共建OPC孵化载体，楚云负责方案、运营、招商、培训、活动、数据和项目孵化。",
            "投资合作：围绕跨境店群、AI内容项目、算力应用项目设立投资、托管、分润和项目库机制。",
            "算力合作：联合光谷超算中心等资源，以渠道销售、应用封装、API平台和行业解决方案实现变现。"
        ],
    },
    {
        "title": "阶段性落地重点",
        "subtitle": "先做可展示、可成交、可复制的样板",
        "bullets": [
            "第一阶段：完善官网、企业宣发PPT、招商材料、AI漫剧课程框架、跨境店群合作方案和算力合作介绍。",
            "第二阶段：打造高校实训与OPC园区样板，形成课程、工位、店铺、内容生产、算力调用和数据看板。",
            "第三阶段：引入跨境服务商、内容团队、算力渠道、投资人和企业客户，形成可持续项目库。",
            "第四阶段：推动分中心复制，形成湖北AI数字贸易与OPC全域孵化标杆。"
        ],
    },
    {
        "title": "官网内容落点",
        "subtitle": "把新增业务自然放进官网，而不是堆概念",
        "bullets": [
            "首页业务版图：展示四大业务，用户第一屏即可理解楚云数航能做什么。",
            "关于我们：说明公司定位、资源渠道、光谷超算中心战略合作和四类核心能力。",
            "合作体系：按高校、园区、投资人、算力伙伴、内容团队分别设置合作入口。",
            "投资与合作页：突出项目阶段、资源合作状态、算力与园区样板、可对接的投资项目。"
        ],
    },
    {
        "title": "结束页",
        "subtitle": "楚云数航：以AI赋能内容，以店群连接全球，以园区承载项目，以算力驱动未来",
        "bullets": [
            "一人起步，平台加持；高校共创，园区承载；AI赋能，数航出海。",
            "湖北楚云数航科技有限公司正在构建面向高校、政府、园区、企业和投资伙伴的AI数字贸易新基础设施。"
        ],
    },
]


def update_home():
    path = ROOT / "static/js/app.js"
    source, start, end, data = load_snapshot(path)
    home = data["home"]
    zh = home["home-content"]["zh-CN"]

    zh["description"] = "湖北楚云数航科技有限公司官网，展示AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售，以及光谷超算中心战略合作等核心业务。"
    zh["hero"]["title"] = "楚云数航<br />AI驱动的跨境电商、AI内容与算力园区运营平台"
    zh["hero"]["subtitle"] = "围绕AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售，连接高校人才、政府园区、跨境渠道、光谷超算中心战略合作资源和产业服务商。"
    zh["news"]["slides"] = [
        {
            "image": "news-robot-home-carousel.jpg",
            "url": "investor.html#announcements",
            "chips": ["业务升级", "光谷超算"],
            "mediaWords": ["AI漫剧", "店群", "算力"],
            "title": "楚云数航四大业务版图完成升级",
            "strong": "产业闭环",
            "body": '<span class="news-copy-nowrap">新增AI漫剧制作与培训、跨境店群投资孵化、OPC园区基建运营、AI算力运营销售，并将光谷超算中心战略合作纳入算力资源体系。</span>',
            "button": "查看详情",
        }
    ]
    zh["productsHeading"] = {
        "title": "业务版图",
        "subtitle": "围绕AI内容、跨境店群、OPC园区和AI算力形成可培训、可投资、可运营、可复制的业务矩阵",
    }
    zh["technology"] = {
        "title": "平台能力",
        "subtitle": "以内容生产、店群运营、园区管理和算力销售四类平台能力支撑业务落地",
        "listLabel": "平台能力列表",
        "cards": [
            {
                "title": "AI内容与漫剧制作中台",
                "body": "沉淀剧本、角色、分镜、素材、模板、任务协同和商业发布流程，服务AI漫剧培训、品牌内容制作与账号矩阵运营。",
                "meta": "内容引擎",
                "button": "了解更多",
                "buttonUrl": "about.html",
            },
            {
                "title": "跨境店群运营管理平台",
                "body": "管理多店铺、订单、库存、客服、选品、财务分润和投资项目复盘，让店群孵化有数据、有风控、有收益模型。",
                "meta": "经营系统",
                "button": "了解更多",
                "buttonUrl": "partner.html#solutions",
            },
            {
                "title": "OPC园区数字化运营平台",
                "body": "承接工位、学员、项目、活动、政策申报、招商展示和政府数据看板，支撑园区基建与长期运营。",
                "meta": "园区底座",
                "button": "了解更多",
                "buttonUrl": "partner.html#recruit",
            },
            {
                "title": "AI算力运营与销售平台",
                "body": "结合光谷超算中心战略合作等资源，承接算力套餐、API调用、模型应用、成本核算和渠道销售。",
                "meta": "算力资源",
                "button": "了解更多",
                "buttonUrl": "investor.html#announcements",
            },
        ],
    }

    products = home["products"]["zh-CN"]
    products["sensenova"] = {
        "number": "01",
        "title": "AI漫剧制作与培训",
        "subtitle": "把AI内容生产变成可学习、可交付、可商业化的训练体系",
        "dynamicBackground": True,
        "links": [
            {"title": "AI漫剧课程体系", "kicker": "AI Drama Training", "desc": "覆盖剧本策划、分镜生成、角色一致性、配音字幕、剪辑发布和账号运营。", "tone": "nova", "url": "partner.html#recruit"},
            {"title": "内容制作交付", "kicker": "Content Studio", "desc": "面向企业宣传、品牌短片、虚拟人内容、短剧漫剧IP提供制作与运营服务。", "tone": "video", "url": "partner.html#solutions"},
            {"title": "高校与园区训练营", "kicker": "Camp", "desc": "为高校、园区、MCN和创业团队输出AI内容创作实训与项目孵化。", "tone": "commerce", "url": "about.html#contact"},
        ],
    }
    products["infra"] = {
        "number": "02",
        "title": "跨境店群投资与孵化",
        "subtitle": "用真实店铺、真实订单和真实数据筛选创业项目",
        "dynamicBackground": True,
        "links": [
            {"title": "店群投资与托管", "kicker": "Store Portfolio", "desc": "围绕TikTok、Ozon、速卖通、独立站等赛道组织店群投资、托管和分润。", "tone": "commerce", "url": "partner.html#solutions"},
            {"title": "运营陪跑与风控", "kicker": "Operation", "desc": "统一选品、刊登、订单、客服、复盘和风控标准，让项目过程可跟踪。", "tone": "nova", "url": "partner.html#benefits"},
            {"title": "项目库与融资撮合", "kicker": "Capital Ready", "desc": "以GMV、利润、复购和成长档案筛选优质店群项目，对接投资人与企业资源。", "tone": "infra", "url": "investor.html#investor-contact"},
        ],
    }
    products["city"] = {
        "number": "03",
        "title": "OPC园区基建与运营",
        "subtitle": "为政府园区和产业载体建设可招商、可孵化、可展示的运营样板",
        "dynamicBackground": True,
        "links": [
            {"title": "OPC空间规划", "kicker": "Park Build", "desc": "规划OPC工位、AI算力实训室、跨境店群运营区、直播间和校企洽谈区。", "tone": "office", "url": "partner.html#solutions"},
            {"title": "园区运营服务", "kicker": "Park Ops", "desc": "输出招商、活动、培训、路演、政策申报、项目台账和数据看板运营。", "tone": "commerce", "url": "investor.html#announcements"},
            {"title": "分中心复制", "kicker": "Replication", "desc": "围绕高校、产业园区和地方政府合作，形成可复制的OPC分中心机制。", "tone": "infra", "url": "partner.html#recruit"},
        ],
    }
    products["innovation"] = {
        "number": "04",
        "title": "AI算力运营及销售",
        "subtitle": "依托光谷超算中心战略合作等资源，连接算力供给与AI应用场景",
        "dynamicBackground": True,
        "links": [
            {"title": "算力资源销售", "kicker": "Compute Sales", "desc": "开展算力套餐、API中转、模型调用服务和企业算力采购对接。", "tone": "infra", "url": "investor.html#announcements"},
            {"title": "AI应用封装", "kicker": "AI Apps", "desc": "围绕AI漫剧、数字人直播、跨境客服、企业知识库和虚拟商品做应用套壳。", "tone": "video", "url": "about.html"},
            {"title": "光谷超算合作", "kicker": "Optics Valley", "desc": "将光谷超算中心战略合作纳入资源背书，支撑算力运营、项目孵化和企业应用落地。", "tone": "nova", "url": "partner.html#benefits"},
        ],
    }
    save_snapshot(path, source, start, end, data)


def update_subpages():
    path = ROOT / "static/js/subpage.js"
    source, start, end, data = load_snapshot(path)
    sub = data["subpage"]
    about = sub["about"]["zh-CN"]
    partner = sub["partner"]["zh-CN"]
    investor = sub["investor"]["zh-CN"]

    about["title"] = "AI驱动的跨境电商、AI内容与算力园区运营平台"
    about["subtitle"] = "楚云数航以AI漫剧制作、跨境店群、OPC园区和AI算力为四大抓手，连接高校、政府园区、光谷超算中心战略合作资源与产业服务商。"
    about["paragraphs"] = [
        "湖北楚云数航科技有限公司定位为AI驱动的跨境电商、AI内容与算力园区运营平台，面向高校、政府园区、创业者、投资人、跨境商家、AI内容团队和企业客户提供综合服务。",
        "公司的四大业务版块包括AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售，形成从学习、制作、开店、运营、算力使用到项目孵化的闭环。",
        "在资源层面，楚云数航将光谷超算中心战略合作纳入算力资源体系，结合API平台、模型应用、企业智能化场景和园区实训空间，支撑AI内容与跨境电商业务落地。",
        "公司希望让学生、创业者和企业用一套平台获得课程、项目、店铺、算力、导师、资金和园区载体支持，让AI产业应用和跨境电商创业从一个人开始，由一套平台放大。",
    ]
    about["stats"] = [["4 大", "业务版块"], ["4 类", "平台能力"], ["光谷超算", "战略合作资源"], ["店群投资", "跨境孵化模型"], ["OPC园区", "基建运营样板"], ["AI漫剧", "制作培训体系"]]
    about["contactSubtitle"] = "欢迎交流AI漫剧培训、跨境店群投资、OPC园区运营、AI算力合作、高校共建与政企资源协同需求。"
    about["officesSubtitle"] = "围绕高校、园区、跨境渠道、AI内容团队、算力资源和产业服务商形成协同网络。"
    about["offices"] = [
        ["武汉光谷", "中国", "公司落地与OPC园区运营核心区域"],
        ["光谷超算中心", "中国", "AI算力运营、销售与企业应用战略合作资源"],
        ["高校基地", "中国", "AI漫剧、数字贸易、跨境店群与SYB双创实训"],
        ["跨境渠道", "全球", "TikTok / Ozon / 速卖通 / 独立站等赛道协同"],
        ["AI内容生态", "中国", "短剧漫剧、虚拟人、数字人直播与内容电商协同"],
    ]

    partner["heroTitle"] = "与楚云数航共建AI内容、跨境店群、OPC园区与算力生态"
    partner["heroSubtitle"] = "面向高校、政府园区、内容团队、跨境商家、投资人与算力伙伴，提供可培训、可投资、可运营、可复制的产业合作方案。"
    partner["why"] = [
        ["01", "四大业务清晰，合作入口明确", "AI漫剧制作培训、跨境店群投资孵化、OPC园区基建运营、AI算力运营销售均可独立合作，也可组合落地。"],
        ["02", "真实项目驱动，成果可量化", "以内容作品、店铺订单、园区活动、算力订单、学员成长档案和项目收益作为核心成果指标。"],
        ["03", "光谷超算资源加持", "将光谷超算中心战略合作纳入算力资源体系，支撑AI漫剧、跨境电商智能化和企业AI应用落地。"],
        ["04", "政校企园可复制样板", "围绕高校课程、园区载体、投资项目、算力销售和产业展示，形成可复制分中心模式。"],
    ]
    partner["solutions"] = [
        ["AI内容", "AI漫剧制作与培训", "联合高校、园区、内容团队开设AI漫剧、虚拟人、短视频、品牌内容和内容电商训练营。", "static/picture/chuyun-brand-placeholder.svg", "AI漫剧制作与培训"],
        ["跨境店群", "店群投资与项目孵化", "面向投资人、企业和创业者提供跨境店群托管、运营陪跑、分润结算和项目库筛选。", "static/picture/chuyun-brand-placeholder.svg", "跨境店群投资与孵化"],
        ["OPC园区", "园区基建与运营", "为政府园区和产业载体规划OPC工位、直播空间、产品展示、AI实训室、路演活动和数据看板。", "static/picture/chuyun-brand-placeholder.svg", "OPC园区基建与运营"],
        ["AI算力", "算力运营及销售", "依托光谷超算中心战略合作等资源，开展算力套餐、API中转、模型应用封装和企业AI解决方案。", "static/picture/chuyun-brand-placeholder.svg", "AI算力运营及销售"],
    ]
    partner["partners"] = [
        ["高校与学院", "共建AI漫剧、数字贸易、跨境店群、SYB/双创课程和实习实训基地。", "提交高校合作意向"],
        ["政府与产业园区", "共建OPC园区载体、AI算力实训室、跨境产品展示窗口、创业赛事和政策落地样板。", "提交园区合作需求"],
        ["投资人与企业客户", "参与跨境店群投资、品牌出海、货源合作、AI内容项目、算力应用项目和收益分成。", "提交投资合作需求"],
        ["算力与AI服务伙伴", "联合光谷超算中心等资源伙伴，共建算力销售、API平台、模型应用和行业解决方案。", "申请算力生态合作"],
    ]
    partner["benefits"] = [
        ["顶层方案共创", "共同梳理AI内容、跨境店群、OPC园区、算力销售的资源边界、收益路径和阶段目标。"],
        ["系统与数据支撑", "接入AI内容中台、店群运营管理、OPC园区数字化平台和AI算力销售管理能力。"],
        ["运营交付陪跑", "围绕课程、招生、内容制作、店铺运营、园区招商、算力销售和项目复盘提供协同。"],
        ["光谷超算与产业资源对接", "根据项目阶段对接光谷超算中心战略合作资源、高校、政府、园区、跨境渠道、资金和产业服务商。"],
    ]
    partner["solutionsSubtitle"] = "从AI漫剧、跨境店群、OPC园区到AI算力，围绕可落地项目输出标准化合作方案"
    partner["recruitSubtitle"] = "高校、园区、投资人、算力伙伴和AI内容团队均可直接进入合作沟通"
    partner["benefitsSubtitle"] = "把合作拆解为课程、项目、空间、算力、数据和收益机制，确保每一步都有交付、有沉淀"
    partner["ctaSubtitle"] = "无论是AI漫剧培训、跨境店群投资、OPC园区运营，还是AI算力销售合作，都可以从一套清晰的产业孵化方案开始。"

    investor["heroEvents"] = [
        {"tag": "业务升级", "date": "2026.07", "title": "AI漫剧、跨境店群、OPC园区、AI算力四大业务版图完成梳理", "image": "static/picture/chuyun-brand-placeholder.svg"},
        {"tag": "战略合作", "date": "推进中", "title": "光谷超算中心战略合作纳入AI算力运营与销售资源体系", "image": "static/picture/chuyun-brand-placeholder.svg"},
        {"tag": "协同方向", "date": "持续开放", "title": "面向高校、园区、投资人、企业客户和算力伙伴开放合作入口", "image": "static/picture/chuyun-brand-placeholder.svg"},
    ]
    investor["governanceDesc"] = "楚云数航当前围绕AI漫剧制作培训、跨境店群投资孵化、OPC园区基建运营、AI算力运营销售和光谷超算中心战略合作资源整合推进。"
    investor["reportsDesc"] = "从品牌官网与宣发材料开始，逐步落地AI漫剧课程、跨境店群项目、OPC园区样板、算力销售通道和光谷超算中心战略合作应用场景。"
    investor["announcementsDesc"] = "围绕AI漫剧、跨境店群、OPC园区、AI算力、光谷超算战略合作和高校产教融合记录阶段进展。"
    investor["announcements"] = [
        ["楚云数航新增AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售四大业务版图", "cn.html#products"],
        ["光谷超算中心战略合作资源纳入AI算力运营与销售体系", "partner.html#benefits"],
        ["企业宣发PPT内容稿与PPTX草稿已同步新增业务与合作方信息", "content-research/chuyun-corporate-deck.md"],
    ]
    investor["stockExchange"] = "项目阶段：四大业务版图与光谷超算战略合作资源整合中 / 开放高校、园区、投资、算力与企业合作"
    investor["stockStaticChange"] = "开放四大业务合作洽谈"
    investor["addresses"] = ["欢迎通过官网联系表单预约沟通AI漫剧、跨境店群、OPC园区和AI算力合作", "可根据区域、园区、算力资源和项目需求安排线上或线下对接"]

    save_snapshot(path, source, start, end, data)


def write_markdown():
    markdown = "# 湖北楚云数航科技有限公司企业宣发PPT内容稿\n\n"
    markdown += "本稿已补齐AI漫剧制作与培训、跨境店群投资与孵化、OPC园区基建与运营、AI算力运营及销售，并将光谷超算中心战略合作纳入资源与合作方体系。\n\n"
    for index, slide in enumerate(slides, 1):
        markdown += f"## {index}. {slide['title']}\n"
        markdown += f"{slide['subtitle']}\n\n"
        for bullet in slide["bullets"]:
            markdown += f"- {bullet}\n"
        markdown += "\n"
    (OUT_DIR / "chuyun-corporate-deck.md").write_text(markdown, encoding="utf-8")


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=(17, 24, 39), align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    return box


def write_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for index, item in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(247, 249, 252)

        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(209, 43, 36)
        accent.line.fill.background()

        add_textbox(slide, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.62), item["title"], size=30, bold=True)
        add_textbox(slide, Inches(0.78), Inches(1.24), Inches(11.6), Inches(0.45), item["subtitle"], size=16, color=(209, 43, 36))

        top = 2.05
        for bullet in item["bullets"]:
            marker = slide.shapes.add_shape(9, Inches(0.86), Inches(top + 0.06), Inches(0.12), Inches(0.12))
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(209, 43, 36)
            marker.line.fill.background()
            add_textbox(slide, Inches(1.08), Inches(top), Inches(11.35), Inches(0.58), bullet, size=15, color=(31, 41, 55))
            top += 0.82

        add_textbox(slide, Inches(11.55), Inches(6.92), Inches(1.0), Inches(0.25), f"{index}/{len(slides)}", size=9, color=(102, 112, 133), align=PP_ALIGN.RIGHT)
    prs.save(OUT_DIR / "chuyun-corporate-deck.pptx")


update_home()
update_subpages()
write_markdown()
write_pptx()
